# pip install python-pptx -i https://pypi.tuna.tsinghua.edu.cn/simple
from pptx import Presentation
import os


def readpptx(path, resultdir):
    content = []
    ppt = Presentation(path)
    imgnum = 0

    for slide in ppt.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                if hasattr(shape, 'image'):
                    imgnum += 1
                    with open(os.path.join(resultdir, str(imgnum) + '-' + shape.image.filename), 'wb') as f:
                        f.write(shape.image.blob)
                        f.close()
            else:
                content.append(shape.text)

    content = '\n'.join(content)
    with open(os.path.join(resultdir, 'text.txt'), 'w') as f:
        f.write(content)
        f.close()
    return content


if __name__ == '__main__':
    path = r'D:/2020增补-2019彭阳年终.pptx'
    resultdir = r'D:/transppt'
    txt = readpptx(path, resultdir)
    print(txt)
