import os
from pptx import Presentation


def readnotes(path):
    content = []
    ppt = Presentation(path)

    for indx, slide in enumerate(ppt.slides):
        try:
            nt = slide.notes_slide.notes_text_frame.text
            content.append((indx, nt))
        except:
            content.append((indx, ''))

    # content = '\n'.join(content)
    return content


def readtxt(path):
    content = []
    ppt = Presentation(path)

    for slide in ppt.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                pass
            else:
                content.append(shape.text)

    # content = '\n'.join(content)
    return content


def readimg(path, savedir, save_prefix=''):
    if not savedir:
        return
    ppt = Presentation(path)
    imgnum = 0
    imginfo = []

    coordinate_scale = 10000  # prevent out of float

    last_text = ''
    for indx, slide in enumerate(ppt.slides):
        # if indx == 41:
        #     print()
        txtinfo_slide = []
        imginfo_slide = []
        curtext = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                if hasattr(shape, 'image'):
                    imgnum += 1
                    fname = save_prefix + '-' + str(indx) + '-' + str(imgnum) + '-' + shape.image.filename
                    imginfo_slide.append(
                        {'fname': fname, 'relatedtxt': '',
                         'x': shape.left / coordinate_scale + shape.width / coordinate_scale / 2,
                         'y': shape.top / coordinate_scale + shape.height / coordinate_scale / 2})

                    # save
                    with open(os.path.join(savedir, fname), 'wb') as f:
                        f.write(shape.image.blob)
                        f.close()
            else:
                curtext.append(shape.text)
                if len(shape.text_frame.paragraphs[0].runs) < 1:
                    continue
                fontsize = 18.0
                if shape.text_frame.paragraphs[0].runs[0].font.size:
                    fontsize = shape.text_frame.paragraphs[0].runs[0].font.size.pt

                shape_left = 0  # for very weird errors
                shape_top = 0
                shape_height = 0
                shape_width = 0
                if shape.left:
                    shape_left = shape.left
                if shape.top:
                    shape_top = shape.top
                if shape.height:
                    shape_height = shape.height
                if shape.width:
                    shape_width = shape.width
                txtinfo_slide.append(
                    {'text': shape.text.replace('\n', ''), 'ftsize': fontsize,
                     'x': shape_left / coordinate_scale + shape_width / coordinate_scale / 2,
                     'y': shape_top / coordinate_scale + shape_height / coordinate_scale / 2})

        # find best related text
        min_related_text = 5
        for cimg in imginfo_slide:
            if not txtinfo_slide:  # no text then give latest page text
                cimg['relatedtxt'] = last_text
                continue
            sorted_txt = sorted(txtinfo_slide,
                                key=lambda t: distance_score(t['x'], t['y'], cimg['x'], cimg['y'], t['ftsize']))
            if len(sorted_txt) < 2 or len(sorted_txt[0]['text']) > min_related_text:
                cimg['relatedtxt'] = sorted_txt[0]['text']
            else:
                cimg['relatedtxt'] = sorted_txt[0]['text'] + ',' + sorted_txt[1]['text']

        curtext = '\n'.join(curtext)
        if len(curtext) > 1:
            last_text = curtext

        # add info in this slide to final result
        imginfo += imginfo_slide
    return imginfo


def distance_score(tx, ty, imgx, imgy, tfontsize):
    return abs(tx - imgx) + abs(ty - imgy)  # /tfontsize


# not in use, for experiment
def readpptx(path, resultdir):
    content = []
    ppt = Presentation(path)
    imgnum = 0

    for slide in ppt.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                if hasattr(shape, 'image'):
                    imgnum += 1
                    with open(os.path.join(resultdir, str(imgnum) + '-' + shape.image.filename), 'wb') as f:
                        f.write(shape.image.blob)
                        f.close()
            else:
                content.append(shape.text)

    content = '\n'.join(content)
    with open(os.path.join(resultdir, 'text.txt'), 'w') as f:
        f.write(content)
        f.close()
    return content
