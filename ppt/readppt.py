from pptx import Presentation



def readpptx(path, extend_table=False, extend_image=False):
    content = []
    ppt = Presentation(path)

    for slide in ppt.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                if extend_image and hasattr(shape, 'image'):
                    with open(shape.image.filename, 'wb') as f:
                        f.write(shape.image.blob)
                        f.close()
                    # 调用图片文字识别
                    # content += readImage(shape.image.filename)
                    # 移除临时图片
                    # os.remove(shape.image.filename)
                # 提取表格内容
                if extend_table and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            content += cell.text
            else:
                content.append(shape.text)

    content = '\n'.join(content)
    return content